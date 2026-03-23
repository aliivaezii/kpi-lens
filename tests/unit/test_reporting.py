"""Unit tests for the Excel and PowerPoint report exporters."""

from __future__ import annotations

from io import BytesIO

import pytest
from openpyxl import load_workbook

from kpi_lens.reporting.excel_exporter import generate_workbook
from kpi_lens.reporting.powerpoint import generate_presentation

# ── Shared fixtures ────────────────────────────────────────────────────────────

SAMPLE_SNAPSHOT = {
    "otif": {
        "value": 96.5,
        "unit": "%",
        "health_status": "green",
        "benchmark_distance": 1.0,
        "display_name": "OTIF Delivery Rate",
        "direction": "higher_is_better",
    },
    "supplier_dppm": {
        "value": 2800.0,
        "unit": "ppm",
        "health_status": "red",
        "benchmark_distance": -250.0,
        "display_name": "Supplier DPPM",
        "direction": "lower_is_better",
    },
}

SAMPLE_ANOMALIES = [
    {
        "id": 1,
        "kpi_name": "supplier_dppm",
        "period_start": "2025-01-06",
        "period_end": "2025-01-12",
        "observed_value": 2800.0,
        "expected_low": 400.0,
        "expected_high": 900.0,
        "severity": 0.85,
        "detector_name": "ensemble",
        "entity": "global",
        "is_acknowledged": False,
    }
]

SAMPLE_TRENDS = {
    "otif": [
        {"period_end": "2025-01-12", "value": 96.5},
        {"period_end": "2025-01-19", "value": 95.1},
    ]
}


# ── Excel exporter ─────────────────────────────────────────────────────────────


def test_generate_workbook_returns_bytes() -> None:
    result = generate_workbook(SAMPLE_SNAPSHOT, SAMPLE_ANOMALIES, SAMPLE_TRENDS)
    assert isinstance(result, bytes)
    assert len(result) > 0


def test_generate_workbook_has_three_sheets() -> None:
    raw = generate_workbook(SAMPLE_SNAPSHOT, SAMPLE_ANOMALIES, SAMPLE_TRENDS)
    wb = load_workbook(BytesIO(raw))
    assert len(wb.sheetnames) == 3


def test_generate_workbook_sheet_names() -> None:
    raw = generate_workbook(SAMPLE_SNAPSHOT, SAMPLE_ANOMALIES, SAMPLE_TRENDS)
    wb = load_workbook(BytesIO(raw))
    assert wb.sheetnames == ["Executive Summary", "Anomaly Detail", "KPI Trends"]


def test_summary_sheet_contains_kpi_names() -> None:
    raw = generate_workbook(SAMPLE_SNAPSHOT, SAMPLE_ANOMALIES, SAMPLE_TRENDS)
    wb = load_workbook(BytesIO(raw))
    ws = wb["Executive Summary"]
    cell_values = [ws.cell(row=r, column=1).value for r in range(1, ws.max_row + 1)]
    assert "otif" in cell_values
    assert "supplier_dppm" in cell_values


def test_anomaly_sheet_contains_severity() -> None:
    raw = generate_workbook(SAMPLE_SNAPSHOT, SAMPLE_ANOMALIES, SAMPLE_TRENDS)
    wb = load_workbook(BytesIO(raw))
    ws = wb["Anomaly Detail"]
    # Row 2 is the first data row (row 1 is header)
    severity_value = ws.cell(row=2, column=8).value
    assert severity_value == pytest.approx(0.85, abs=0.01)


def test_generate_workbook_empty_anomalies() -> None:
    """Exporter must not raise when anomaly list is empty."""
    result = generate_workbook(SAMPLE_SNAPSHOT, [], SAMPLE_TRENDS)
    wb = load_workbook(BytesIO(result))
    ws = wb["Anomaly Detail"]
    # Only the header row should exist
    assert ws.max_row == 1


def test_generate_workbook_empty_snapshot() -> None:
    """Exporter must handle an empty snapshot gracefully."""
    result = generate_workbook({}, [], {})
    wb = load_workbook(BytesIO(result))
    assert len(wb.sheetnames) == 3


# ── PowerPoint exporter ────────────────────────────────────────────────────────


def test_generate_presentation_returns_bytes() -> None:
    result = generate_presentation(SAMPLE_SNAPSHOT, SAMPLE_ANOMALIES)
    assert isinstance(result, bytes)
    assert len(result) > 0


def test_generate_presentation_has_three_slides() -> None:
    from pptx import Presentation

    raw = generate_presentation(SAMPLE_SNAPSHOT, SAMPLE_ANOMALIES)
    prs = Presentation(BytesIO(raw))
    assert len(prs.slides) == 3


def test_generate_presentation_no_anomalies() -> None:
    """Presentation with zero anomalies must not raise."""
    raw = generate_presentation(SAMPLE_SNAPSHOT, [])
    from pptx import Presentation

    prs = Presentation(BytesIO(raw))
    assert len(prs.slides) == 3


def test_generate_presentation_custom_date() -> None:
    """report_date parameter should be accepted without error."""
    raw = generate_presentation(
        SAMPLE_SNAPSHOT, SAMPLE_ANOMALIES, report_date="2025-03-23"
    )
    assert isinstance(raw, bytes)
