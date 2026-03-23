"""
PowerPoint report generator — produces an executive-ready slide deck.

Design decisions:
- Uses python-pptx's blank layout (index 6) to avoid theme conflicts when
  the file is opened on machines with different Office versions installed.
- Text boxes are positioned in EMU (English Metric Units = 914400 per inch)
  so slide geometry is explicit and portable across python-pptx versions.
- One anomaly per bullet point on the Anomaly slide; truncated to 10 items
  to keep the deck readable in a board-room setting.
"""

from __future__ import annotations

from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from kpi_lens.kpis.definitions import (
    KPI_BY_NAME,  # noqa: F401 — reserved for future per-KPI slide pages
)

# Slide dimensions for the default 16:9 widescreen layout (in EMU)
_SLIDE_WIDTH = Emu(9144000)
_SLIDE_HEIGHT = Emu(5143500)

# Brand colours matching the Streamlit dashboard palette
_COLOUR_PRIMARY = RGBColor(0x2F, 0x4F, 0x8F)  # Navy — headers and accents
_COLOUR_GREEN = RGBColor(0x70, 0xAD, 0x47)  # Status green
_COLOUR_YELLOW = RGBColor(0xFF, 0xC0, 0x00)  # Status yellow
_COLOUR_RED = RGBColor(0xFF, 0x00, 0x00)  # Status red
_COLOUR_BODY = RGBColor(0x26, 0x26, 0x26)  # Near-black for body text

_STATUS_COLOUR: dict[str, RGBColor] = {
    "green": _COLOUR_GREEN,
    "yellow": _COLOUR_YELLOW,
    "red": _COLOUR_RED,
}


def _add_text_box(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    bold: bool = False,
    colour: RGBColor | None = None,
    align: Any = PP_ALIGN.LEFT,
) -> None:
    tx_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tx_box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = colour or _COLOUR_BODY


def _add_title_slide(prs: Presentation, report_date: str, anomaly_count: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Background rectangle
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0),
        Inches(0),
        _SLIDE_WIDTH,
        _SLIDE_HEIGHT,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = _COLOUR_PRIMARY
    bg.line.fill.background()

    _add_text_box(
        slide,
        "KPI-Lens Supply Chain Report",
        0.5,
        1.8,
        9.0,
        1.0,
        font_size=36,
        bold=True,
        colour=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide,
        f"Period ending {report_date}  ·  {anomaly_count} active anomalies",
        0.5,
        2.9,
        9.0,
        0.6,
        font_size=18,
        colour=RGBColor(0xCC, 0xDD, 0xFF),
        align=PP_ALIGN.CENTER,
    )


def _add_kpi_summary_slide(prs: Presentation, snapshot: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(
        slide,
        "KPI Health Summary",
        0.3,
        0.2,
        9.4,
        0.6,
        font_size=24,
        bold=True,
        colour=_COLOUR_PRIMARY,
    )

    # Table: KPI | Value | Unit | Status
    row_count = len(snapshot) + 1  # +1 for header
    table = slide.shapes.add_table(
        row_count,
        4,
        Inches(0.3),
        Inches(1.0),
        Inches(9.4),
        Inches(0.45 * row_count),
    ).table

    for col_idx, header in enumerate(["KPI", "Value", "Unit", "Status"]):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
            0xFF, 0xFF, 0xFF
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = _COLOUR_PRIMARY

    for row_idx, (kpi_name, data) in enumerate(snapshot.items(), start=1):
        status = data.get("health_status", "")
        row_data = [
            data.get("display_name", kpi_name),
            str(round(data.get("value", 0), 2))
            if data.get("value") is not None
            else "—",
            data.get("unit", ""),
            status.upper() if status else "—",
        ]
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            if col_idx == 3 and status in _STATUS_COLOUR:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _STATUS_COLOUR[status]


def _add_anomaly_slide(prs: Presentation, anomalies: list[dict[str, Any]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(
        slide,
        "Active Anomalies",
        0.3,
        0.2,
        9.4,
        0.6,
        font_size=24,
        bold=True,
        colour=_COLOUR_PRIMARY,
    )

    # Show at most 10 anomalies — beyond that a board-room audience loses the thread
    displayed = anomalies[:10]
    bullet_lines = []
    for anomaly in displayed:
        kpi = anomaly.get("kpi_name", "").upper()
        severity_pct = int(anomaly.get("severity", 0) * 100)
        period = anomaly.get("period_start", "")
        observed = anomaly.get("observed_value", "—")
        bullet_lines.append(
            f"• {kpi} | Severity {severity_pct}% | {period} | Observed: {observed}"
        )

    if not bullet_lines:
        bullet_lines = ["• No anomalies detected in the selected period."]

    _add_text_box(
        slide,
        "\n".join(bullet_lines),
        0.5,
        1.1,
        9.0,
        3.8,
        font_size=14,
        colour=_COLOUR_BODY,
    )


def generate_presentation(
    snapshot: dict[str, Any],
    anomalies: list[dict[str, Any]],
    report_date: str | None = None,
) -> bytes:
    """
    Build a 3-slide executive deck and return its bytes.

    All data comes from the repository layer — this function is a pure
    transformation from dicts to pptx bytes, with no I/O of its own.
    """
    from datetime import date as date_type

    effective_date = report_date or str(date_type.today())
    prs = Presentation()
    prs.slide_width = _SLIDE_WIDTH
    prs.slide_height = _SLIDE_HEIGHT

    _add_title_slide(prs, effective_date, len(anomalies))
    _add_kpi_summary_slide(prs, snapshot)
    _add_anomaly_slide(prs, anomalies)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
